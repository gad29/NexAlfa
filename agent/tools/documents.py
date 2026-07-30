"""
NexAlfa Document Tools
Read, write, convert, and extract from PDF, DOCX, TXT, CSV, XLSX, JSON, YAML, HTML, MD, PPTX.
"""

from __future__ import annotations

import json
import logging
import os
from pathlib import Path
from typing import Optional

from agent.tools.base import Tool

logger = logging.getLogger("nex.tools.documents")


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _read_pdf(path: Path) -> str:
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"--- Page {i+1} ---\n{page_text}")
        return "\n\n".join(text_parts) if text_parts else "(No readable text found in PDF)"
    except ImportError:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(path))
            parts = []
            for i, page in enumerate(reader.pages):
                t = page.extract_text() or ""
                if t.strip():
                    parts.append(f"--- Page {i+1} ---\n{t}")
            return "\n\n".join(parts) if parts else "(No readable text found)"
        except ImportError:
            return "ERROR: Install pdfplumber or PyPDF2: pip install pdfplumber"


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else ""
                if "Heading" in style:
                    level = style.replace("Heading ", "").strip()
                    prefix = "#" * int(level) if level.isdigit() else "##"
                    parts.append(f"{prefix} {para.text}")
                else:
                    parts.append(para.text)
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n| " + " | ".join(["---"] * len(table.columns)) + " |")
                for r in rows:
                    parts.append(f"| {r} |")
        return "\n\n".join(parts) if parts else "(Empty document)"
    except ImportError:
        return "ERROR: Install python-docx: pip install python-docx"


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cells))
            parts.append("\n".join(rows[:500]))
            if ws.max_row and ws.max_row > 500:
                parts.append(f"... ({ws.max_row - 500} more rows)")
        wb.close()
        return "\n\n".join(parts)
    except ImportError:
        return "ERROR: Install openpyxl: pip install openpyxl"


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
            if texts:
                parts.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
        return "\n\n".join(parts) if parts else "(Empty presentation)"
    except ImportError:
        return "ERROR: Install python-pptx: pip install python-pptx"


def _read_csv(path: Path) -> str:
    import csv
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = []
        for i, row in enumerate(reader):
            if i > 500:
                rows.append(f"... (truncated, {i}+ rows total)")
                break
            rows.append(" | ".join(row))
    return "\n".join(rows)


READERS = {
    ".txt": _read_text, ".md": _read_text, ".markdown": _read_text,
    ".json": _read_text, ".yaml": _read_text, ".yml": _read_text,
    ".html": _read_text, ".htm": _read_text, ".xml": _read_text,
    ".csv": _read_csv, ".tsv": _read_csv,
    ".pdf": _read_pdf,
    ".docx": _read_docx, ".doc": _read_docx,
    ".xlsx": _read_xlsx, ".xls": _read_xlsx,
    ".pptx": _read_pptx, ".ppt": _read_pptx,
    ".py": _read_text, ".js": _read_text, ".ts": _read_text,
    ".java": _read_text, ".c": _read_text, ".cpp": _read_text,
    ".rs": _read_text, ".go": _read_text, ".rb": _read_text,
    ".sh": _read_text, ".bat": _read_text, ".ps1": _read_text,
    ".ini": _read_text, ".cfg": _read_text, ".toml": _read_text,
    ".env": _read_text, ".log": _read_text,
}

SUPPORTED_FORMATS = ", ".join(sorted(set(READERS.keys())))


# ── Tool Classes ───────────────────────────────────────────

class DocReadTool(Tool):
    name = "doc_read"
    description = "Read any document file. Supports: PDF, DOCX, TXT, CSV, XLSX, PPTX, JSON, YAML, MD, HTML, and code files."

    def get_schema(self) -> dict:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to the document file"},
                    "max_chars": {"type": "integer", "description": "Max characters to return (default 50000)"},
                },
                "required": ["path"],
            },
        }

    async def execute(self, path: str, max_chars: int = 50000) -> str:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return f"ERROR: File not found: {path}"
        ext = p.suffix.lower()
        reader = READERS.get(ext)
        if not reader:
            return f"ERROR: Unsupported format '{ext}'. Supported: {SUPPORTED_FORMATS}"
        try:
            content = reader(p)
            if len(content) > max_chars:
                content = content[:max_chars] + f"\n\n... (truncated at {max_chars} chars, total {len(content)})"
            return f"📄 **{p.name}** ({ext}, {p.stat().st_size:,} bytes)\n\n{content}"
        except Exception as e:
            return f"ERROR: Failed to read {path}: {type(e).__name__}: {e}"


class DocWriteTool(Tool):
    name = "doc_write"
    description = "Create/write a document. Format auto-detected from extension. Supports: TXT, MD, JSON, CSV, HTML, DOCX, PDF, XLSX. Use markdown-style headings (# ## ###) for DOCX/PDF."

    def get_schema(self) -> dict:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Output file path (extension determines format)"},
                    "content": {"type": "string", "description": "The document content to write"},
                    "format": {"type": "string", "description": "Force format (default: auto-detect from extension)"},
                },
                "required": ["path", "content"],
            },
        }

    async def execute(self, path: str, content: str, format: str = "auto") -> str:
        p = Path(path).expanduser().resolve()
        p.parent.mkdir(parents=True, exist_ok=True)
        ext = format if format != "auto" else p.suffix.lower().lstrip(".")
        try:
            if ext in ("txt", "md", "markdown", "html", "htm", "csv", "json", "yaml", "yml",
                        "xml", "py", "js", "ts", "sh", "bat", "ini", "cfg", "toml", "log", "env"):
                p.write_text(content, encoding="utf-8")
            elif ext == "docx":
                from docx import Document
                doc = Document()
                for line in content.split("\n"):
                    line = line.strip()
                    if line.startswith("# "):
                        doc.add_heading(line[2:], level=1)
                    elif line.startswith("## "):
                        doc.add_heading(line[3:], level=2)
                    elif line.startswith("### "):
                        doc.add_heading(line[4:], level=3)
                    elif line:
                        doc.add_paragraph(line)
                doc.save(str(p))
            elif ext == "pdf":
                from reportlab.lib.pagesizes import A4
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                doc = SimpleDocTemplate(str(p), pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                for line in content.split("\n"):
                    line = line.strip()
                    if line.startswith("# "):
                        story.append(Paragraph(line[2:], styles["Heading1"]))
                    elif line.startswith("## "):
                        story.append(Paragraph(line[3:], styles["Heading2"]))
                    elif line:
                        story.append(Paragraph(line, styles["Normal"]))
                    story.append(Spacer(1, 6))
                doc.build(story)
            elif ext == "xlsx":
                from openpyxl import Workbook
                import csv, io
                wb = Workbook()
                ws = wb.active
                reader = csv.reader(io.StringIO(content))
                for row in reader:
                    ws.append(row)
                wb.save(str(p))
            else:
                p.write_text(content, encoding="utf-8")
            size = p.stat().st_size
            return f"✅ Written: {p} ({size:,} bytes)"
        except Exception as e:
            return f"ERROR: Failed to write {path}: {type(e).__name__}: {e}"


class DocConvertTool(Tool):
    name = "doc_convert"
    description = "Convert a document between formats. Example: convert report.docx to report.pdf"

    def get_schema(self) -> dict:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "source": {"type": "string", "description": "Source file path"},
                    "target": {"type": "string", "description": "Target file path (extension determines output format)"},
                },
                "required": ["source", "target"],
            },
        }

    async def execute(self, source: str, target: str) -> str:
        src = Path(source).expanduser().resolve()
        if not src.exists():
            return f"ERROR: Source file not found: {source}"
        read_tool = DocReadTool()
        content = await read_tool.execute(str(src), max_chars=500000)
        if content.startswith("ERROR:"):
            return content
        lines = content.split("\n", 2)
        raw_content = lines[2] if len(lines) > 2 else content
        write_tool = DocWriteTool()
        return await write_tool.execute(target, raw_content)


class DocExtractTablesTool(Tool):
    name = "doc_extract_tables"
    description = "Extract all tables from a PDF or DOCX file as structured text."

    def get_schema(self) -> dict:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PDF or DOCX file"},
                },
                "required": ["path"],
            },
        }

    async def execute(self, path: str) -> str:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return f"ERROR: File not found: {path}"
        ext = p.suffix.lower()
        tables_found = []
        if ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(str(p)) as pdf:
                    for i, page in enumerate(pdf.pages):
                        for j, table in enumerate(page.extract_tables()):
                            rows = [" | ".join(str(c or "") for c in row) for row in table]
                            tables_found.append(f"### Table {j+1} (Page {i+1})\n" + "\n".join(rows))
            except ImportError:
                return "ERROR: Install pdfplumber: pip install pdfplumber"
        elif ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(str(p))
                for j, table in enumerate(doc.tables):
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    tables_found.append(f"### Table {j+1}\n" + "\n".join(rows))
            except ImportError:
                return "ERROR: Install python-docx: pip install python-docx"
        else:
            return f"Table extraction not supported for {ext}. Use PDF or DOCX."
        if not tables_found:
            return "No tables found in the document."
        return f"📊 Found {len(tables_found)} table(s):\n\n" + "\n\n".join(tables_found)


def get_document_tools() -> list[Tool]:
    return [DocReadTool(), DocWriteTool(), DocConvertTool(), DocExtractTablesTool()]
